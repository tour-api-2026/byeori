"""out.pptx 를 그 자리에서 고친다. 이미지가 이미 들어가 있어 재생성하면 잃는다.
서식을 지키려고 run.text 만 갈아끼운다(text_frame.text 대입은 서식이 날아간다)."""
import copy, sys
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.text.text import _Paragraph

BODY = RGBColor(0x1A, 0x1A, 0x1F)

def style(run, size):
    run.font.color.rgb = BODY
    run.font.bold = False
    if size: run.font.size = Pt(size)

def set_cell(cell, text, size=None):
    tf = cell.text_frame
    lines = text.split('\n')
    p0 = tf.paragraphs[0]
    if not p0.runs: p0.add_run()
    p0.runs[0].text = lines[0]; style(p0.runs[0], size)
    for extra in p0.runs[1:]: extra._r.getparent().remove(extra._r)
    for p in tf.paragraphs[1:]: p._p.getparent().remove(p._p)
    for line in lines[1:]:
        newp = copy.deepcopy(p0._p)
        p0._p.getparent().append(newp)
        para = _Paragraph(newp, tf)
        para.runs[0].text = line; style(para.runs[0], size)

def table_of(slide, rows, cols):
    for sh in slide.shapes:
        if sh.has_table and len(sh.table.rows) == rows and len(sh.table.columns) == cols:
            return sh.table
    sys.exit(f'표를 찾지 못함 {rows}x{cols}')

prs = Presentation('work2.pptx')
S = prs.slides

# ── 3쪽 핵심기능: 건수 갱신 + 지도/검색이 실시간임을 명시 ──
set_cell(table_of(S[2], 1, 2).cell(0, 1),
         '- 전통 테마 행사 큐레이션 기능\n'
         '   공연·축제 데이터에서 전통을 주제로 한 행사만 자동 분류해 1,340건 제공\n'
         '- 위치 기반 지도 탐색 기능\n'
         '   지도가 멈춘 좌표·반경으로 공사 OpenAPI를 실시간 조회해 그 영역의 장소만 표시\n'
         '- 명칭 키워드 검색 기능\n'
         '   입력한 검색어로 공사 OpenAPI를 실시간 조회. 숙박·쇼핑·레포츠까지 전 유형 대상\n'
         '- 장소 상세 정보 제공 기능\n'
         '   소개글·이용시간·휴무일·문의처를 실시간 조회하고, 진행 중인 행사를 함께 제공\n'
         '- 여행 일정 생성 및 경로 안내 기능\n'
         '   방문지를 담아 하루 일정을 구성하고 실제 이동 경로와 소요 시간 확인\n'
         '- 리뷰 및 이용자 보호 기능\n'
         '   방문 후기 작성, 부적절한 콘텐츠 신고, 특정 이용자 차단')

# ── 10쪽 API 목록: 실시간 3건을 앞에 두고 정기 수집 2건을 뒤에 ──
APIS = [
    ('한국관광공사 국문 관광정보 서비스 (locationBasedList2)',
     '지도 주변 탐색에 실시간 호출. 지도가 멈춘 중심 좌표와 반경으로 그 영역의 장소만 조회. 지도 이동 1회당 1회 호출'),
    ('한국관광공사 국문 관광정보 서비스 (searchKeyword2)',
     '검색 화면에 실시간 호출. 이용자가 입력한 검색어로 조회하며, 입력이 멈춘 뒤 0.4초 후 1회 호출'),
    ('한국관광공사 국문 관광정보 서비스 (detailCommon2 / detailIntro2)',
     '장소 상세 화면을 열 때마다 실시간 호출하여 소개글·이용시간·휴무일·문의처·주차 정보를 조회. 장소 1건 조회당 2회 호출'),
    ('한국관광공사 국문 관광정보 서비스 (areaBasedList2)',
     '지역 기반 관광지·문화시설·음식점 정보를 1일 1회 수집. 관심사 기반 추천과 지역별 탐색에 활용 (29,537건 반영)'),
    ('한국관광공사 국문 관광정보 서비스 (searchFestival2 / ldongCode2)',
     '행사·축제 정보와 법정동 코드를 1일 1회 수집. 전통 테마 행사 분류와 지역 매핑에 활용 (417건 반영)'),
]
t = table_of(S[9], 10, 3)
for i, (name, desc) in enumerate(APIS):
    set_cell(t.cell(i * 2, 2), name, size=12)
    set_cell(t.cell(i * 2 + 1, 2), desc, size=12)

# ── 12쪽 차별성: 3번을 실제 구조에 맞게 ──
set_cell(table_of(S[11], 2, 2).cell(0, 1),
         '1. 장소 정보와 공연 정보의 결합\n'
         '   대부분의 관광 서비스는 장소 또는 공연 한쪽만 다룹니다. 벼리는 한국관광공사 OpenAPI의 장소\n'
         '   데이터와 공연 데이터를 연결해, 장소 상세에서 "여기서 지금 무엇이 열리는지"를 함께 보여줍니다.\n'
         '2. 전통 행사 자동 분류\n'
         '   공공데이터는 전통 행사를 따로 구분해 주지 않습니다. 장르 코드만으로는 놓치는 행사가 많아\n'
         '   제목·주최기관을 함께 검사하는 규칙으로 1,340건을 선별했습니다.\n'
         '3. 필요한 만큼만 부르는 실시간 조회\n'
         '   지도·검색·상세를 모두 공사 OpenAPI 실시간 호출로 구현했습니다. 전국 데이터를 미리 받아두지\n'
         '   않고 보고 있는 영역만 부르며, 좌표를 반올림해 같은 화면을 다시 부르지 않습니다.\n'
         '4. 한복 착용 혜택 정보\n'
         '   한복 착용 시 혜택이 있는 장소를 별도로 표시해 전통문화 체험이 하나의 동선으로 이어지게 했습니다.\n'
         '5. 진입 장벽 없는 이용\n'
         '   장소·공연 조회와 지도 탐색은 로그인 없이 모두 이용할 수 있습니다.\n'
         '6. 이용자 위치 정보 비수집\n'
         '   현재 위치 기능은 기기 안에서만 좌표를 사용하고 서버로 전송하지 않습니다.', size=10.5)

# ── 6쪽 흐름도 설명(기능 2)도 실시간으로 ──
head = table_of(S[5], 2, 2)
set_cell(head.cell(0, 1), '위치 기반 지도 탐색', size=14)
set_cell(head.cell(1, 1),
         '지도가 멈춘 좌표와 반경으로 공사 OpenAPI를 실시간 조회해 그 영역의 장소만 표시합니다. '
         '전국 데이터를 미리 받아두지 않으며, 위치 정보는 지도 이동에만 쓰고 서버로 전송하지 않습니다.', size=10.5)
body = table_of(S[5], 3, 4)
set_cell(body.cell(2, 0), '위치 권한 요청 → 지도 표시 → 실시간 주변 조회 → 마커 표시 → 상세 이동', size=12)

prs.save('edited.pptx')
print('텍스트 수정 완료 → edited.pptx')

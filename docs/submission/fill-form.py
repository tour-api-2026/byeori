"""기능설명서 양식(작성용).pptx 에 벼리 내용을 채운다.

양식의 서식을 지키기 위해:
- 셀의 첫 run 에만 텍스트를 넣고 나머지 run 은 지운다(run.text 대입 — text_frame.text 는 서식이 날아간다)
- 여러 줄은 첫 문단을 복제해 <a:pPr> 를 물려받게 한다
- 안내용 회색 박스는 도형째 제거한다
"""
import copy
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

DECK = 'work.pptx'
OUT = 'out.pptx'
IMG = '/home/hidi/dev/byeori/docs'


BODY = RGBColor(0x1A, 0x1A, 0x1F)


def style(run, size):
    """가이드 문구의 빨간 서식이 그대로 상속되므로 본문 서식을 명시한다."""
    run.font.color.rgb = BODY
    run.font.bold = False
    if size:
        run.font.size = Pt(size)


def set_cell(cell, text, size=None):
    """셀 텍스트 교체. 줄바꿈은 문단으로 나눈다."""
    tf = cell.text_frame
    lines = text.split('\n')
    p0 = tf.paragraphs[0]
    runs = p0.runs
    if not runs:
        p0.add_run()
        runs = p0.runs
    runs[0].text = lines[0]
    style(runs[0], size)
    for extra in runs[1:]:
        extra._r.getparent().remove(extra._r)
    for p in tf.paragraphs[1:]:
        p._p.getparent().remove(p._p)
    for line in lines[1:]:
        newp = copy.deepcopy(p0._p)
        p0._p.getparent().append(newp)
        from pptx.text.text import _Paragraph
        para = _Paragraph(newp, tf)
        para.runs[0].text = line
        style(para.runs[0], size)


def drop_shape(slide, name_starts):
    for sh in list(slide.shapes):
        if sh.name.startswith(name_starts):
            sh._element.getparent().remove(sh._element)
            return True
    return False


def table_of(slide, rows, cols):
    for sh in slide.shapes:
        if sh.has_table and len(sh.table.rows) == rows and len(sh.table.columns) == cols:
            return sh.table
    raise SystemExit(f'표를 찾지 못함 {rows}x{cols}')


prs = Presentation(DECK)
S = prs.slides

# ── 1. 표지 ──────────────────────────────────────────────
t = table_of(S[0], 2, 2)
set_cell(t.cell(0, 1), '벼리')          # 팀명 — 콘텐츠랩 등록값으로 교체 필요
set_cell(t.cell(1, 1), '벼리')          # 서비스명

# ── 2. 서비스 소개 ① ──────────────────────────────────────
t = table_of(S[1], 4, 2)
set_cell(t.cell(0, 1), '벼리')
set_cell(t.cell(1, 1), '웹 서비스')
set_cell(t.cell(2, 1),
         '전국의 관광 명소와 전통 공연·축제를 지도에서 함께 찾아보고 '
         '하루 여행 일정까지 만드는 한국 전통문화 여행 가이드 서비스')
set_cell(t.cell(3, 1),
         '전통 공연·축제 정보와 관광 명소 정보는 서로 다른 기관이 각각 제공해 흩어져 있습니다. '
         '국악 공연을 보러 가면서 근처 고궁이나 한옥 카페를 함께 둘러보려 해도 여러 사이트를 '
         '오가며 위치를 대조해야 합니다.\n'
         '전통문화를 접하기 번거로운 이유는 콘텐츠가 부족해서가 아니라 정보가 연결되어 있지 '
         '않기 때문이라고 보았습니다. 한국관광공사 OpenAPI에는 전국 관광 정보가, '
         '공연예술통합전산망에는 공연 정보가 이미 충분히 공개되어 있습니다.\n'
         "벼리는 이 둘을 하나의 지도 위에 결합해 전통문화를 '보러 가는 일'이 아니라 "
         "'여행하는 일'로 만들고자 했습니다. '벼리'는 그물을 오므렸다 폈다 하는 굵은 줄을 뜻하는 "
         '순우리말로, 흩어진 정보를 하나로 엮는다는 뜻을 담았습니다.')

# ── 3. 서비스 핵심기능 ────────────────────────────────────
t = table_of(S[2], 1, 2)
set_cell(t.cell(0, 1),
         '- 전통 테마 행사 큐레이션 기능\n'
         '   공연·축제 데이터에서 전통을 주제로 한 행사만 자동 분류해 1,340건 제공\n'
         '- 위치 기반 지도 탐색 기능\n'
         '   현재 위치를 중심으로 주변 관광지·문화시설·공연장·음식점을 지도에 표시\n'
         '- 장소 상세 정보 제공 기능\n'
         '   한국관광공사 OpenAPI를 실시간 조회해 소개글·이용시간·휴무일·문의처를 표시하고,\n'
         '   해당 장소에서 진행 중인 행사를 함께 제공\n'
         '- 여행 일정 생성 및 경로 안내 기능\n'
         '   방문지를 담아 하루 일정을 구성하고 실제 이동 경로와 소요 시간 확인\n'
         '- 리뷰 및 이용자 보호 기능\n'
         '   방문 후기 작성, 부적절한 콘텐츠 신고, 특정 이용자 차단')

# ── 4. 서비스 이미지 (텍스트만 비우고 그림은 뒤에서 넣는다) ──
t = table_of(S[3], 2, 2)
set_cell(t.cell(0, 1), '')
set_cell(t.cell(1, 1), '')

# ── 5~9. 핵심 기능별 흐름도 ───────────────────────────────
FLOWS = [
    ('전통 테마 행사 큐레이션',
     '공연·축제 데이터에서 전통 주제 행사만 자동 분류해 제공하는 기능입니다. 장르 코드에 더해 제목·주최기관을 함께 검사해 1,340건을 선별했습니다.',
     '01-전통행사-큐레이션.png',
     '공연·축제 수집 → 전통 여부 판별 → 전통 행사 저장 → 목록 제공 → 사용자 탐색'),
    ('위치 기반 지도 탐색',
     '현재 위치 중심으로 주변 관광지·문화시설·공연장·음식점을 지도에 표시합니다. 위치 정보는 지도 이동에만 쓰고 서버로 전송하지 않습니다.',
     '02-지도-주변탐색.png',
     '위치 권한 요청 → 지도 표시 → 장소 조회 → 마커 표시 → 상세 이동'),
    ('장소 상세 정보 제공',
     '화면을 열 때 한국관광공사 OpenAPI를 실시간 조회해 소개글·이용시간·휴무일·문의처를 최신 값으로 보여주고, 진행 중인 행사도 함께 표시합니다.',
     '03-장소-상세정보.png',
     '장소 선택 → 공사 OpenAPI 실시간 조회 → 연관 행사 결합 → 이용자 정보 결합 → 화면 구성'),
    ('여행 일정 생성 및 경로 안내',
     '가고 싶은 곳을 담아 하루 일정을 구성합니다. 방문 순서를 정하면 실제 이동 경로를 계산해 지도에 그리고 총 거리·소요 시간을 보여줍니다.',
     '04-여행일정-경로.png',
     '일정 생성 → 방문지 추가 → 순서 정리 → 경로 탐색 → 지도 확인'),
    ('리뷰 및 이용자 보호',
     '방문 후기를 별점과 함께 남깁니다. 부적절한 리뷰·장소는 신고할 수 있고, 이용자를 차단하면 그 사람의 리뷰가 보이지 않습니다.',
     '05-리뷰-신고차단.png',
     '리뷰 작성 → 평점 재계산 → 신고 → 차단 → 차단 관리'),
]
for i, (title, desc, _img, steps) in enumerate(FLOWS):
    sl = S[4 + i]
    drop_shape(sl, '직사각형')                  # 회색 안내 박스 제거
    head = table_of(sl, 2, 2)
    set_cell(head.cell(0, 0), f'핵심 기능{i + 1}')
    set_cell(head.cell(0, 1), title, size=14)
    set_cell(head.cell(1, 1), desc, size=10.5)
    body = table_of(sl, 3, 4)
    # 흐름도 이미지와 단계 설명은 4칸을 가로질러 쓴다. 병합하지 않으면 빈 칸 3개가 남는다.
    body.cell(1, 0).merge(body.cell(1, 3))
    body.cell(2, 0).merge(body.cell(2, 3))
    set_cell(body.cell(1, 0), '')              # 이미지 자리 — 뒤에서 그림 삽입
    set_cell(body.cell(2, 0), steps, size=12)

# ── 10. 한국관광공사 OpenAPI ──────────────────────────────
t = table_of(S[9], 10, 3)
APIS = [
    ('한국관광공사 국문 관광정보 서비스 (detailCommon2)',
     '장소 상세 화면을 열 때마다 실시간 호출하여 소개글·홈페이지 정보를 조회. 장소 1건 조회당 1회 호출'),
    ('한국관광공사 국문 관광정보 서비스 (detailIntro2)',
     '장소 상세 화면을 열 때마다 실시간 호출하여 이용시간·휴무일·문의처·주차 정보를 조회. 장소 1건 조회당 1회 호출'),
    ('한국관광공사 국문 관광정보 서비스 (areaBasedList2)',
     '지역 기반 관광지·문화시설·음식점 정보 수집. 전국 단위 지도 탐색과 카테고리 필터에 활용 (29,535건 반영)'),
    ('한국관광공사 국문 관광정보 서비스 (searchFestival2)',
     '행사·축제 정보 수집. 전통 테마 행사 분류와 장소별 진행 중인 행사 표시에 활용'),
    ('한국관광공사 국문 관광정보 서비스 (ldongCode2 / lDongSignguCd)',
     '법정동·시군구 코드 조회. 주소를 지역으로 매핑하여 지역별 탐색 필터에 활용'),
]
for i, (name, desc) in enumerate(APIS):
    set_cell(t.cell(i * 2, 2), name, size=12)
    set_cell(t.cell(i * 2 + 1, 2), desc, size=12)

# ── 11. 기타 데이터 ───────────────────────────────────────
sl = S[10]
drop_shape(sl, '직사각형')
t = table_of(sl, 6, 3)
ETC = [
    ('공연예술통합전산망(KOPIS) 공연 정보',
     '(재)예술경영지원센터 제공. 공연 정보 7,953건을 수집해 전통 공연 분류의 기반으로 활용'),
    ('서울시 문화행사 정보',
     '서울시 열린데이터광장 제공. 문화행사 737건을 수집해 행사 정보 보강에 활용'),
    ('카카오맵 / 카카오모빌리티 API',
     '지도 표시와 장소 키워드 검색, 여행 일정의 이동 경로·소요 시간 계산에 활용'),
]
for i, (name, desc) in enumerate(ETC):
    set_cell(t.cell(i * 2, 2), name, size=12)
    set_cell(t.cell(i * 2 + 1, 2), desc, size=12)

# ── 12. 차별성 & 발전계획 ─────────────────────────────────
t = table_of(S[11], 2, 2)
set_cell(t.cell(0, 1),
         '1. 장소 정보와 공연 정보의 결합\n'
         '   대부분의 관광 서비스는 장소 또는 공연 한쪽만 다룹니다. 벼리는 한국관광공사 OpenAPI의 장소\n'
         '   데이터와 공연 데이터를 연결해, 장소 상세에서 "여기서 지금 무엇이 열리는지"를 함께 보여줍니다.\n'
         '2. 전통 행사 자동 분류\n'
         '   공공데이터는 전통 행사를 따로 구분해 주지 않습니다. 장르 코드만으로는 놓치는 행사가 많아\n'
         '   제목·주최기관을 함께 검사하는 규칙으로 1,340건을 선별했습니다.\n'
         '3. 실시간 조회와 정기 수집의 병행\n'
         '   전국 지도 탐색은 정기 수집 데이터로 빠르게 제공하고, 개별 장소의 운영 정보는 상세 화면을\n'
         '   열 때 공사 OpenAPI를 실시간 조회해 최신 값을 보여줍니다. 속도와 최신성을 함께 확보했습니다.\n'
         '4. 한복 착용 혜택 정보\n'
         '   한복 착용 시 혜택이 있는 장소를 별도로 표시해 전통문화 체험이 하나의 동선으로 이어지게 했습니다.\n'
         '5. 진입 장벽 없는 이용\n'
         '   장소·공연 조회와 지도 탐색은 로그인 없이 모두 이용할 수 있습니다.\n'
         '6. 이용자 위치 정보 비수집\n'
         '   현재 위치 기능은 기기 안에서만 좌표를 사용하고 서버로 전송하지 않습니다.', size=10.5)
set_cell(t.cell(1, 1),
         '1. 다국어 지원\n'
         '   한국관광공사 OpenAPI는 영문·일문·중문·독일어 등을 동일 인증키로 제공합니다. 호출 대상만\n'
         '   교체하면 되므로 방한 외국인 관광객 대상 다국어 서비스로 확장합니다.\n'
         '2. 지역별 심화 큐레이션\n'
         '   전주 한옥마을 일대, 안동 하회마을 권역 등 지역 단위 전통문화 테마 코스로 심화해\n'
         '   지역 관광 활성화에 기여합니다.\n'
         '3. 무장애 여행 정보 결합\n'
         '   무장애 여행 API(KorWithService2)를 결합해 고령자·장애인도 전통문화를 편히 즐길 수 있도록\n'
         '   접근성 정보를 제공합니다.\n'
         '4. 개인화 추천\n'
         '   즐겨찾기와 방문 이력을 기반으로 관심 지역·유형에 맞는 코스를 추천합니다.\n'
         '5. 모바일 앱 배포\n'
         '   React Native 기반으로 안드로이드 앱 빌드가 완료되었으며, Google Play 정식 등록을 진행합니다.', size=10.5)

prs.save(OUT)
print('저장:', OUT)
